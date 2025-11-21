#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation about the Jakarta SC 2025 Website
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
import os

def create_jakarta_sc_presentation():
    """Create a comprehensive PowerPoint presentation about Jakarta SC 2025 Website"""
    
    # Create a new presentation
    prs = Presentation()
    
    # Define color scheme
    primary_blue = RGBColor(33, 150, 243)  # Material Blue
    dark_blue = RGBColor(21, 101, 192)
    light_gray = RGBColor(245, 245, 245)
    dark_gray = RGBColor(66, 66, 66)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "Jakarta SC 2025"
    title1.text_frame.paragraphs[0].font.size = Pt(44)
    title1.text_frame.paragraphs[0].font.color.rgb = primary_blue
    title1.text_frame.paragraphs[0].font.bold = True
    
    subtitle1.text = "Public Facilities Finder Web Application\nPresentation Overview"
    subtitle1.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle1.text_frame.paragraphs[0].font.color.rgb = dark_gray
    
    # Slide 2: Agenda
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content2 = slide2.placeholders[1]
    
    title2.text = "Agenda Presentasi"
    title2.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    agenda_items = [
        "• Gambaran Umum Proyek",
        "• Fitur dan Fungsionalitas",
        "• Teknologi yang Digunakan",
        "• Struktur Aplikasi",
        "• Kategori Fasilitas",
        "• Lokalisasi Multi-Bahasa",
        "• Responsivitas Mobile",
        "• QR Code Integration",
        "• Demo dan Kesimpulan"
    ]
    
    content2.text = "\n".join(agenda_items)
    for paragraph in content2.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
    
    # Slide 3: Project Overview
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    content3 = slide3.placeholders[1]
    
    title3.text = "Gambaran Umum Proyek"
    title3.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    overview_text = """• Nama Aplikasi: Jakarta SC 2025 - Public Facilities Finder
• Tujuan: Membantu pengunjung ICE BSD menemukan fasilitas umum terdekat
• Platform: Progressive Web Application (PWA)
• Target Users: Peserta dan pengunjung acara di ICE BSD
• Lokasi: ICE BSD International Convention Exhibition
• Alamat: Jl. BSD Grand Boulevard No.1, Pagedangan, Tangerang"""
    
    content3.text = overview_text
    for paragraph in content3.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
    
    # Slide 4: Key Features
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    content4 = slide4.placeholders[1]
    
    title4.text = "Fitur Utama Aplikasi"
    title4.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    features_text = """✓ Pencarian Fasilitas Berdasarkan Kategori
✓ Informasi Lokasi dengan Google Maps Integration
✓ Estimasi Jarak dan Waktu Tempuh
✓ Support 14 Bahasa Internasional
✓ Responsive Design untuk Mobile & Desktop
✓ QR Code Generator untuk Sharing
✓ Navigasi Intuitif dengan Material-UI
✓ Progressive Web App (PWA) Features"""
    
    content4.text = features_text
    for paragraph in content4.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = dark_gray
    
    # Slide 5: Technology Stack
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    content5 = slide5.placeholders[1]
    
    title5.text = "Teknologi yang Digunakan"
    title5.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    tech_text = """Frontend Framework:
• React 19.1.0 dengan TypeScript
• Vite sebagai Build Tool
• React Router DOM untuk Navigation

UI/UX Libraries:
• Material-UI (MUI) v7.2.0
• Tailwind CSS v3.4.1
• Material Icons

Additional Libraries:
• QRCode generation library
• Canvas API untuk QR customization
• React Hooks untuk State Management"""
    
    content5.text = tech_text
    for paragraph in content5.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
    
    # Slide 6: Facility Categories
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    content6 = slide6.placeholders[1]
    
    title6.text = "Kategori Fasilitas yang Tersedia"
    title6.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    categories_text = """🏧 ATM - Lokasi mesin ATM terdekat
🏥 Hospital - Rumah sakit dan klinik kesehatan  
💊 Pharmacy - Apotek dan toko obat
🍽️ Restaurant - Restoran dan tempat makan
⛽ Gas Station - SPBU dan pom bensin
💱 Money Changer - Tempat penukaran mata uang
🔧 Auto Repair - Bengkel dan service mobil

Setiap kategori dilengkapi dengan:
• Informasi provider/penyedia layanan
• Detail lokasi dengan alamat lengkap
• Koordinat GPS untuk navigasi
• Estimasi jarak dan waktu tempuh"""
    
    content6.text = categories_text
    for paragraph in content6.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
    
    # Slide 7: Internationalization
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    content7 = slide7.placeholders[1]
    
    title7.text = "Dukungan Multi-Bahasa (Internationalization)"
    title7.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    languages_text = """Aplikasi mendukung 14 bahasa internasional:

🇬🇧 English (Default)          🇯🇵 日本語 (Japanese)
🇮🇩 Bahasa Indonesia         🇲🇾 Bahasa Malaysia  
🇲🇲 မြန်မာဘာသာ (Myanmar)      🇳🇱 Nederlands (Dutch)
🇵🇭 Filipino                  🇹🇼 繁體中文 (Traditional Chinese)
🇹🇭 ไทย (Thai)                🇧🇷 Português (Brasil)
🇨🇴 Español (Colombia)        🇮🇳 हिन्दी (Hindi)
🇮🇹 Italiano (Italian)         🇱🇰 සිංහල (Sinhala)

Fitur Lokalisasi:
• Dynamic language switching
• URL-based locale routing (/en, /id, /ja, etc.)
• Localized content untuk semua UI elements"""
    
    content7.text = languages_text
    for paragraph in content7.text_frame.paragraphs:
        paragraph.font.size = Pt(12)
    
    # Slide 8: Mobile Responsiveness
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    content8 = slide8.placeholders[1]
    
    title8.text = "Mobile-First Responsive Design"
    title8.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    responsive_text = """Mobile Optimization Features:
• Adaptive card layouts untuk berbagai screen sizes
• Touch-friendly interface dengan minimum 44px touch targets
• Responsive typography dan spacing
• Optimized untuk mobile performance

Desktop Enhancements:
• Larger information cards dengan better spacing
• Enhanced hover effects dan animations  
• Multi-column layouts untuk better content organization
• Desktop-specific navigation features

Cross-Platform Compatibility:
• Progressive Web App (PWA) capabilities
• Offline-ready functionality
• App-like experience di mobile devices
• Fast loading dengan optimized assets"""
    
    content8.text = responsive_text
    for paragraph in content8.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
    
    # Slide 9: QR Code Integration
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    content9 = slide9.placeholders[1]
    
    title9.text = "QR Code Integration"
    title9.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    qr_text = """QR Code Generator Features:
• Custom QR codes dengan logo embedding
• Downloadable QR codes dalam format PNG
• Branded QR codes dengan Jakarta SC logo
• Easy sharing untuk promotional materials

Technical Implementation:
• QRCode library untuk generation
• HTML5 Canvas untuk custom rendering
• Logo overlay dengan proper positioning
• Error correction level optimization

Use Cases:
• Marketing materials dan promotional content
• Event signage dan banners  
• Social media sharing
• Print materials untuk event"""
    
    content9.text = qr_text
    for paragraph in content9.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
    
    # Slide 10: Application Architecture
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    content10 = slide10.placeholders[1]
    
    title10.text = "Struktur Aplikasi dan Arsitektur"
    title10.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    architecture_text = """Component Structure:
• HomePage: Landing page dengan kategori fasilitas
• CategoryPage: List providers dalam kategori tertentu  
• ProviderPage: Detail lokasi dengan maps integration
• Header: Navigation dengan language switcher
• QRCodeWithLogo: QR code generator component

Data Management:
• JSON-based data structure untuk setiap kategori
• Centralized data dalam facilitiesData object
• Modular organization dengan separate files
• Easy maintenance dan updates

Routing System:
• React Router DOM untuk navigation
• Locale-based routing (/:locale/:category/:provider)
• Dynamic route generation
• SEO-friendly URL structure"""
    
    content10.text = architecture_text
    for paragraph in content10.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
    
    # Slide 11: Key Benefits
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    content11 = slide11.placeholders[1]
    
    title11.text = "Manfaat dan Keunggulan"
    title11.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    benefits_text = """Untuk Event Organizer:
✓ Meningkatkan experience pengunjung event
✓ Reduce inquiries about facility locations
✓ Professional digital solution
✓ Easy marketing dengan QR codes

Untuk Pengunjung:
✓ Quick access ke informasi fasilitas terdekat
✓ Multilingual support untuk international guests  
✓ Mobile-friendly untuk akses on-the-go
✓ Accurate directions dengan Google Maps

Technical Advantages:
✓ Modern web technologies untuk performance optimal
✓ Scalable architecture untuk future enhancements
✓ SEO-optimized untuk better discoverability
✓ Cross-platform compatibility"""
    
    content11.text = benefits_text
    for paragraph in content11.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
    
    # Slide 12: Future Enhancements
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    title12 = slide12.shapes.title
    content12 = slide12.placeholders[1]
    
    title12.text = "Rencana Pengembangan Futue"
    title12.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    future_text = """Potential Enhancements:
• Real-time availability updates untuk facilities
• User reviews dan ratings system
• Push notifications untuk special announcements  
• Integration dengan booking systems
• Advanced filtering dan search capabilities
• Offline maps untuk areas dengan poor connectivity

Additional Features:
• Event schedule integration
• Emergency contact information
• Accessibility information for disabled users
• Integration dengan transportation apps
• Analytics dashboard untuk usage tracking
• Multi-tenant support untuk different events"""
    
    content12.text = future_text
    for paragraph in content12.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
    
    # Slide 13: Conclusion
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    title13 = slide13.shapes.title
    content13 = slide13.placeholders[1]
    
    title13.text = "Kesimpulan"
    title13.text_frame.paragraphs[0].font.color.rgb = primary_blue
    
    conclusion_text = """Jakarta SC 2025 Public Facilities Finder adalah solusi digital 
modern yang dirancang untuk meningkatkan experience pengunjung 
ICE BSD Convention Center.

Key Highlights:
• Modern React-based web application
• Support 14 bahasa internasional  
• Mobile-first responsive design
• 7 kategori fasilitas lengkap
• QR code integration untuk easy sharing
• Google Maps integration untuk navigation

Aplikasi ini ready untuk deployment dan dapat diakses melalui:
https://jakartasc2025.info

Terima kasih atas perhatiannya!"""
    
    content13.text = conclusion_text
    for paragraph in content13.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
    
    # Slide 14: Q&A
    slide14 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add large centered text for Q&A
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    qa_textbox = slide14.shapes.add_textbox(left, top, width, height)
    qa_text_frame = qa_textbox.text_frame
    qa_text_frame.text = "Questions & Answers"
    
    qa_paragraph = qa_text_frame.paragraphs[0]
    qa_paragraph.font.size = Pt(48)
    qa_paragraph.font.bold = True
    qa_paragraph.font.color.rgb = primary_blue
    qa_paragraph.alignment = PP_ALIGN.CENTER
    
    qa_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add subtitle
    subtitle_box = slide14.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Diskusi dan Pertanyaan"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = dark_gray
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    output_path = "Jakarta_SC_2025_Presentation.pptx"
    prs.save(output_path)
    
    return output_path

if __name__ == "__main__":
    try:
        # Create the presentation
        output_file = create_jakarta_sc_presentation()
        print(f"✅ PowerPoint presentation berhasil dibuat: {output_file}")
        print(f"📁 File location: {os.path.abspath(output_file)}")
        print("\n📋 Presentation Contents:")
        print("1. Title Slide")
        print("2. Agenda")
        print("3. Project Overview") 
        print("4. Key Features")
        print("5. Technology Stack")
        print("6. Facility Categories")
        print("7. Internationalization")
        print("8. Mobile Responsiveness")
        print("9. QR Code Integration")
        print("10. Application Architecture")
        print("11. Key Benefits")
        print("12. Future Enhancements")
        print("13. Conclusion")
        print("14. Q&A")
        
    except ImportError as e:
        print("❌ Error: python-pptx library not found.")
        print("📦 Please install it using: pip install python-pptx")
        print(f"Error details: {e}")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")